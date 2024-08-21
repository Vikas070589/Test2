#!/usr/bin/env python
# coding: utf-8

from flask import Flask, render_template, request, redirect, url_for, send_from_directory, flash
import pandas as pd
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

app = Flask(__name__)
app.secret_key = 'your_secret_key_here'  # Replace with your secret key

UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'output'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['OUTPUT_FOLDER'] = OUTPUT_FOLDER

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        try:
            # Get the files and form data
            excel_file = request.files['excel_file']
            pptx_template_file = request.files['pptx_template_file']
            start_row = int(request.form['start_row'])
            end_row = int(request.form['end_row']) if request.form['end_row'] else None
            output_folder = request.form['output_folder'] or 'pptx_files'

            # Save the uploaded files
            excel_path = os.path.join(app.config['UPLOAD_FOLDER'], excel_file.filename)
            pptx_template_path = os.path.join(app.config['UPLOAD_FOLDER'], pptx_template_file.filename)
            excel_file.save(excel_path)
            pptx_template_file.save(pptx_template_path)

            # Process the files
            output_folder_path = generate_ppt_from_excel(excel_path, pptx_template_path, start_row, end_row, output_folder)

            flash(f"PPT files generated successfully in {output_folder_path}", "success")
            return redirect(url_for('index'))

        except Exception as e:
            flash(str(e), "danger")
            return redirect(url_for('index'))

    return render_template('index.html')


def generate_ppt_from_excel(excel_file, pptx_template, start_row=0, end_row=None, output_folder='pptx_files'):
    df = pd.read_excel(excel_file, sheet_name='Summaries')
    if not end_row:
        end_row = len(df)

    output_folder_path = os.path.join(app.config['OUTPUT_FOLDER'], output_folder)
    os.makedirs(output_folder_path, exist_ok=True)

    for index, row in df.iloc[start_row:end_row].iterrows():
        prs = Presentation(pptx_template)

        shapes_data = {}
        for column in df.columns:
            if column in row:
                value = row[column]
                if pd.notna(value):
                    if column == 'Duckers Solution':
                        shapes_data[column] = replace_bullet_points(value)
                    else:
                        shapes_data[column] = value
                else:
                    shapes_data[column] = ""

        for slide in prs.slides:
            update_shapes_with_excel_data(slide, shapes_data)

        output_pptx_file = os.path.join(output_folder_path, f"{row.get('Case Study Name', 'Slide')}.pptx")
        prs.save(output_pptx_file)
        print(f"Saved {output_pptx_file} in {output_folder_path}")

    return output_folder_path


def replace_bullet_points(text):
    # Replace '*' with '•'
    return text.replace('*', '• ') if text is not None else ""


def update_shapes_with_excel_data(slide, shapes_data):
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape_name = shape.name
            if shape_name in shapes_data:
                shape.text_frame.text = shapes_data[shape_name] if shapes_data[shape_name] is not None else ""
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(18)
                        paragraph.alignment = PP_ALIGN.LEFT


@app.route('/download/<filename>')
def download_file(filename):
    return send_from_directory(app.config['OUTPUT_FOLDER'], filename)


if __name__ == '__main__':
    app.run(debug=True)
